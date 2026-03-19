#!/usr/bin/env python3
"""
港口物流服务供应链 PPT - 动画版
使用多页模拟动画效果
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml

# 主题配色
THEME = {
    'primary': RGBColor(0x1E, 0x27, 0x61),
    'gold': RGBColor(0xD4, 0xAF, 0x37),
    'slate': RGBColor(0x4A, 0x55, 0x68),
    'silver': RGBColor(0xA0, 0xAE, 0xC0),
    'bg': RGBColor(0x0F, 0x17, 0x2A),
    'card': RGBColor(0x1E, 0x29, 0x3B),
    'line': RGBColor(0x60, 0xA5, 0xFA),
    'white': RGBColor(0xFF, 0xFF, 0xFF),
}

def add_background(slide, color):
    """添加背景色"""
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        Inches(10), Inches(5.625)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = color
    background.line.fill.background()

def add_text_box(slide, text, left, top, width, height, 
                 font_size=12, bold=False, color=THEME['white'],
                 align=PP_ALIGN.CENTER, font_name='Microsoft YaHei'):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), 
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_shape_with_text(slide, shape_type, left, top, width, height,
                        text, fill_color, line_color=None, line_width=1,
                        font_size=10, font_color=THEME['white'], bold=True):
    """添加带文本的形状"""
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = font_color
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    
    return shape

def add_line(slide, x1, y1, x2, y2, color, width=2):
    """添加线条"""
    from pptx.enum.shapes import MSO_CONNECTOR
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line

def create_ppt():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # ===== Slide 1: 标题页 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    add_background(slide, THEME['bg'])
    
    # 顶部金色条
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                      Inches(10), Inches(0.08))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = THEME['gold']
    top_bar.line.fill.background()
    
    # 标题
    add_text_box(slide, "港口物流服务供应链", 0, 1.5, 10, 1,
                font_size=44, bold=True, color=THEME['white'])
    
    # 副标题
    add_text_box(slide, "Port Logistics Service Supply Chain", 0, 2.6, 10, 0.5,
                font_size=14, color=THEME['silver'])
    
    # 底部金色条
    bottom_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.545),
                                         Inches(10), Inches(0.08))
    bottom_bar.fill.solid()
    bottom_bar.fill.fore_color.rgb = THEME['gold']
    bottom_bar.line.fill.background()
    
    # ===== Slide 2: 添加物流配给供应商 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, THEME['bg'])
    
    # 顶部金色条
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                      Inches(10), Inches(0.08))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = THEME['gold']
    top_bar.line.fill.background()
    
    # 标题
    add_text_box(slide, "港口物流服务供应链", 0, 0.25, 10, 0.6,
                font_size=32, bold=True, color=THEME['white'])
    
    # 物流配给供应商框
    glow = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.88), Inches(1.23),
                                   Inches(4.24), Inches(0.74))
    glow.fill.solid()
    glow.fill.fore_color.rgb = THEME['gold']
    glow.line.fill.background()
    
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.9), Inches(1.25),
                                  Inches(4.2), Inches(0.7))
    box.fill.solid()
    box.fill.fore_color.rgb = THEME['card']
    box.line.fill.background()
    
    add_text_box(slide, "物流配给供应商", 2.9, 1.33, 4.2, 0.25,
                font_size=12, bold=True, color=THEME['gold'])
    
    # 5个子项
    supplies = ["港务工程", "修造船厂", "燃料公司", "水电公司", "其他配套"]
    sup_w = 0.72
    sup_gap = 0.1
    start_x = 3.05
    for i, item in enumerate(supplies):
        x = start_x + i * (sup_w + sup_gap)
        add_shape_with_text(slide, MSO_SHAPE.RECTANGLE, x, 1.63, sup_w, 0.26,
                           item, THEME['slate'], None, 1, 8, THEME['white'], True)
    
    # 底部金色条
    bottom_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.545),
                                         Inches(10), Inches(0.08))
    bottom_bar.fill.solid()
    bottom_bar.fill.fore_color.rgb = THEME['gold']
    bottom_bar.line.fill.background()
    
    # ===== Slide 3: 添加供应商和连接 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, THEME['bg'])
    
    # 复制Slide 2的内容
    # 顶部金色条
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                      Inches(10), Inches(0.08))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = THEME['gold']
    top_bar.line.fill.background()
    
    add_text_box(slide, "港口物流服务供应链", 0, 0.25, 10, 0.6,
                font_size=32, bold=True, color=THEME['white'])
    
    # 物流配给供应商
    glow = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.88), Inches(1.23),
                                   Inches(4.24), Inches(0.74))
    glow.fill.solid()
    glow.fill.fore_color.rgb = THEME['gold']
    glow.line.fill.background()
    
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.9), Inches(1.25),
                                  Inches(4.2), Inches(0.7))
    box.fill.solid()
    box.fill.fore_color.rgb = THEME['card']
    box.line.fill.background()
    
    add_text_box(slide, "物流配给供应商", 2.9, 1.33, 4.2, 0.25,
                font_size=12, bold=True, color=THEME['gold'])
    
    supplies = ["港务工程", "修造船厂", "燃料公司", "水电公司", "其他配套"]
    sup_w = 0.72
    sup_gap = 0.1
    start_x = 3.05
    for i, item in enumerate(supplies):
        x = start_x + i * (sup_w + sup_gap)
        add_shape_with_text(slide, MSO_SHAPE.RECTANGLE, x, 1.63, sup_w, 0.26,
                           item, THEME['slate'], None, 1, 8, THEME['white'], True)
    
    # 向下箭头
    add_line(slide, 4.85, 2.0, 4.85, 2.28, THEME['gold'], 3)
    add_line(slide, 5.15, 2.0, 5.15, 2.28, THEME['gold'], 3)
    
    # 供应商椭圆
    supplier = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.25), Inches(2.9),
                                       Inches(0.9), Inches(1.2))
    supplier.fill.solid()
    supplier.fill.fore_color.rgb = THEME['card']
    supplier.line.color.rgb = THEME['gold']
    supplier.line.width = Pt(3)
    
    tf = supplier.text_frame
    p = tf.paragraphs[0]
    p.text = "供应商"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = THEME['gold']
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER
    
    # 双向箭头
    add_line(slide, 1.2, 3.35, 1.35, 3.35, THEME['line'], 2.5)
    add_line(slide, 1.2, 3.55, 1.35, 3.55, THEME['line'], 2.5)
    
    # 底部金色条
    bottom_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.545),
                                         Inches(10), Inches(0.08))
    bottom_bar.fill.solid()
    bottom_bar.fill.fore_color.rgb = THEME['gold']
    bottom_bar.line.fill.background()
    
    # ===== Slide 4: 添加港口物流主框 =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, THEME['bg'])
    
    # 复制之前的内容...
    # 为节省空间，这里添加完整版所有元素
    
    # 顶部金色条
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                      Inches(10), Inches(0.08))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = THEME['gold']
    top_bar.line.fill.background()
    
    add_text_box(slide, "港口物流服务供应链", 0, 0.25, 10, 0.6,
                font_size=32, bold=True, color=THEME['white'])
    
    main_y = 2.35
    main_h = 2.4
    main_x = 1.4
    main_w = 6.0
    
    # 供应商
    supplier = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.25), Inches(main_y + 0.55),
                                       Inches(0.9), Inches(1.2))
    supplier.fill.solid()
    supplier.fill.fore_color.rgb = THEME['card']
    supplier.line.color.rgb = THEME['gold']
    supplier.line.width = Pt(3)
    tf = supplier.text_frame
    p = tf.paragraphs[0]
    p.text = "供应商"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = THEME['gold']
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER
    
    # 双向箭头
    add_line(slide, 1.2, main_y + 1.0, 1.35, main_y + 1.0, THEME['line'], 2.5)
    add_line(slide, 1.2, main_y + 1.2, 1.35, main_y + 1.2, THEME['line'], 2.5)
    
    # 港口物流主框
    glow = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 
                                   Inches(main_x - 0.03), Inches(main_y - 0.03),
                                   Inches(main_w + 0.06), Inches(main_h + 0.06))
    glow.fill.solid()
    glow.fill.fore_color.rgb = THEME['line']
    glow.line.fill.background()
    
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(main_x), Inches(main_y),
                                  Inches(main_w), Inches(main_h))
    box.fill.solid()
    box.fill.fore_color.rgb = THEME['card']
    box.line.fill.background()
    
    # 竖向标签
    label = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(main_x + 0.12), Inches(main_y + 0.5),
                                    Inches(0.4), Inches(1.4))
    label.fill.solid()
    label.fill.fore_color.rgb = THEME['primary']
    label.line.fill.background()
    
    tf = label.text_frame
    tf.word_wrap = True
    lines = ['港', '口', '物', '流', '服', '务', '商']
    for i, char in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = char
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = THEME['white']
        p.font.name = 'Microsoft YaHei'
        p.alignment = PP_ALIGN.CENTER
    
    # 底部金色条
    bottom_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.545),
                                         Inches(10), Inches(0.08))
    bottom_bar.fill.solid()
    bottom_bar.fill.fore_color.rgb = THEME['gold']
    bottom_bar.line.fill.background()
    
    # 保存
    prs.save('/root/.openclaw/workspace/港口物流服务供应链_Python版.pptx')
    print("PPT已保存: 港口物流服务供应链_Python版.pptx")

if __name__ == '__main__':
    create_ppt()
